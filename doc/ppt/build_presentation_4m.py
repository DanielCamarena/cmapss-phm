from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from textwrap import wrap

import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
PROJECT_ROOT = ROOT.parents[1]
ASSET_DIR = ROOT / "assets_4m"
OUTPUT_PPTX = ROOT / "final_presentation_4m.pptx"
OUTPUT_NOTES = ROOT / "speaker_notes_4m.md"

BG = RGBColor(248, 243, 232)
TEXT = RGBColor(28, 36, 52)
MUTED = RGBColor(101, 109, 120)
ACCENT = RGBColor(135, 92, 36)
TEAL = RGBColor(20, 111, 105)
RED = RGBColor(181, 65, 65)
PANEL = RGBColor(255, 255, 255)


@dataclass
class Slide:
    title: str
    subtitle: str
    bullets: list[str]
    visual_type: str
    notes: list[str]
    sources: list[str]


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def load_metrics() -> dict:
    return {
        "global": pd.read_csv(PROJECT_ROOT / "out/predictive_layer/04_metrics_global_by_model.csv"),
        "latency": pd.read_csv(PROJECT_ROOT / "out/predictive_layer/04_latency_by_model.csv"),
    }


def save_plot(fig, path: Path) -> None:
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#F8F3E8")
    plt.close(fig)


def build_problem_diagram(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    ax.axis("off")

    ax.text(
        0.18,
        0.58,
        "Benchmark RUL\nprediction",
        ha="center",
        va="center",
        fontsize=20,
        fontweight="bold",
        color="#1C2434",
        bbox=dict(boxstyle="round,pad=0.6", fc="#FFFFFF", ec="#885C24", lw=2),
    )
    ax.text(
        0.82,
        0.58,
        "PHM deployment\nneeds more",
        ha="center",
        va="center",
        fontsize=20,
        fontweight="bold",
        color="#1C2434",
        bbox=dict(boxstyle="round,pad=0.6", fc="#FFFFFF", ec="#156F69", lw=2),
    )
    ax.annotate("", xy=(0.67, 0.58), xytext=(0.33, 0.58), arrowprops=dict(arrowstyle="->", lw=3, color="#885C24"))

    chips = [
        (0.64, 0.28, "Uncertainty", "#156F69"),
        (0.80, 0.28, "Decision logic", "#156F69"),
        (0.96, 0.28, "Traceability", "#156F69"),
    ]
    for x, y, label, color in chips:
        ax.text(
            x,
            y,
            label,
            ha="center",
            va="center",
            fontsize=12,
            color="#FFFFFF",
            bbox=dict(boxstyle="round,pad=0.3", fc=color, ec=color),
        )

    ax.text(
        0.5,
        0.88,
        "NOT sufficient for PHM deployment",
        ha="center",
        va="center",
        fontsize=24,
        fontweight="bold",
        color="#B54141",
    )
    save_plot(fig, path)


def build_architecture_diagram(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    ax.axis("off")

    blocks = [
        (0.18, "Predictive", "RUL + confidence", "#F6E3C8", "#885C24"),
        (0.50, "Agent", "risk + recommendation", "#D8EFEA", "#156F69"),
        (0.82, "Dashboard", "user + audit", "#F6E3C8", "#885C24"),
    ]
    for x, title, sub, fill, edge in blocks:
        ax.text(
            x,
            0.54,
            f"{title}\n{sub}",
            ha="center",
            va="center",
            fontsize=18,
            fontweight="bold",
            color="#1C2434",
            bbox=dict(boxstyle="round,pad=0.7", fc=fill, ec=edge, lw=2),
        )
    ax.annotate("", xy=(0.38, 0.54), xytext=(0.28, 0.54), arrowprops=dict(arrowstyle="->", lw=3, color="#156F69"))
    ax.annotate("", xy=(0.70, 0.54), xytext=(0.60, 0.54), arrowprops=dict(arrowstyle="->", lw=3, color="#156F69"))
    ax.text(0.5, 0.85, "Predictive \u2192 Agent \u2192 Dashboard", ha="center", fontsize=25, fontweight="bold", color="#1C2434")
    save_plot(fig, path)


def build_model_tradeoff_chart(metrics: dict, path: Path) -> None:
    global_df = metrics["global"].copy()
    latency_df = metrics["latency"].copy()
    plot_df = global_df.merge(latency_df[["model_name", "ms_per_sample"]], on="model_name", how="left")
    fig, ax = plt.subplots(figsize=(8.5, 5))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    colors = {"rf": "#156F69", "gru": "#B54141", "hgb": "#885C24", "lstm": "#8392A5"}
    for _, row in plot_df.iterrows():
        ax.scatter(row["rmse"], row["ms_per_sample"], s=180, color=colors.get(row["model_name"], "#885C24"))
        ax.text(row["rmse"] + 0.08, row["ms_per_sample"] + 0.1, row["model_name"].upper(), fontsize=10, weight="bold")
    ax.set_xlabel("RMSE")
    ax.set_ylabel("ms / sample")
    ax.set_title("Best model \u2260 deployed model")
    ax.set_yscale("log")
    ax.grid(alpha=0.2)
    save_plot(fig, path)


def build_dashboard_collage(path: Path) -> None:
    files = [
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_4.png",
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_5.png",
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_9.png",
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_8.png",
    ]
    labels = ["Summary", "Analysis", "Scenarios", "Audit"]
    images = [Image.open(p).convert("RGB") for p in files]
    target = (700, 380)
    tiles = [ImageOps.fit(img, target, method=Image.Resampling.LANCZOS) for img in images]
    canvas = Image.new("RGB", (target[0] * 2 + 40, target[1] * 2 + 80), (248, 243, 232))
    positions = [(0, 0), (target[0] + 40, 0), (0, target[1] + 80), (target[0] + 40, target[1] + 80)]
    for tile, (x, y), label in zip(tiles, positions, labels):
        canvas.paste(tile, (x, y + 30))
        fig, ax = plt.subplots()
        plt.close(fig)
    # draw labels with pillow basic approach
    try:
        from PIL import ImageDraw, ImageFont

        draw = ImageDraw.Draw(canvas)
        font = ImageFont.load_default()
        for (x, y), label in zip(positions, labels):
            draw.text((x + 10, y + 5), label, fill=(28, 36, 52), font=font)
    except Exception:
        pass
    canvas.save(path)


def build_scenario_chart(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(8.5, 4.8))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    labels = ["Baseline", "Scenario"]
    rul = [123.955, 123.874]
    ax.bar(labels, rul, color=["#885C24", "#156F69"], width=0.55)
    ax.set_ylim(0, 145)
    ax.set_ylabel("RUL")
    ax.set_title("What-if example: \u0394RUL \u2248 0, risk unchanged")
    for i, v in enumerate(rul):
        ax.text(i, v + 2, f"{v:.2f}", ha="center", fontsize=10, weight="bold")
    ax.text(0.5, 25, "Risk Score: 35 \u2192 35", ha="center", fontsize=13, color="#B54141", weight="bold")
    save_plot(fig, path)


def build_assets(metrics: dict) -> dict[str, Path]:
    assets = {
        "problem": ASSET_DIR / "problem_diagram.png",
        "architecture": ASSET_DIR / "architecture_diagram.png",
        "tradeoff": ASSET_DIR / "model_tradeoff.png",
        "dashboard": ASSET_DIR / "dashboard_collage.png",
        "scenario": ASSET_DIR / "scenario_chart.png",
    }
    build_problem_diagram(assets["problem"])
    build_architecture_diagram(assets["architecture"])
    build_model_tradeoff_chart(metrics, assets["tradeoff"])
    build_dashboard_collage(assets["dashboard"])
    build_scenario_chart(assets["scenario"])
    return assets


def slide_specs() -> list[Slide]:
    return [
        Slide(
            "Traceable Predictive Maintenance for NASA C-MAPSS",
            "4-minute overview of the deployed PHM application",
            [
                "Benchmark: NASA C-MAPSS turbofan degradation",
                "Goal: turn RUL prediction into usable PHM decisions",
                "Focus: deployment, traceability, and system integration",
            ],
            "title",
            [
                "This talk presents the project as a complete PHM application, not just a benchmark model.",
                "The key idea is to connect prediction, decision logic, and user-facing traceability in one system.",
            ],
            ["doc/paper/main.tex", "README.md"],
        ),
        Slide(
            "Problem",
            "Why benchmark prediction alone is not enough",
            [
                "Benchmark RUL models produce point predictions.",
                "Real PHM systems also need uncertainty, decision logic, and traceability.",
                "Prediction alone is NOT sufficient for PHM deployment.",
            ],
            "problem",
            [
                "This is the core motivation slide.",
                "A low error score is useful, but it still does not tell an operator how confident the system is, what action to take, or how to audit the decision.",
                "That gap is what this project addresses.",
            ],
            ["doc/ppt/02_slide_content_4m.md", "doc/paper/04_discussion.md"],
        ),
        Slide(
            "Task Context",
            "NASA C-MAPSS defines the prediction setting",
            [
                "Run-to-failure engine trajectories",
                "Multivariate sensor + operational data",
                "Goal: predict Remaining Useful Life",
            ],
            "minimal",
            [
                "C-MAPSS gives us degradation trajectories at the engine-unit level.",
                "Each row combines cycle context, operational settings, and sensors, and the task is to estimate remaining useful life before failure.",
            ],
            ["data/readme.txt", "doc/paper/main.tex"],
        ),
        Slide(
            "System Idea",
            "Not just a model, but a system",
            [
                "Predictive \u2192 Agent \u2192 Dashboard",
                "Predictive: RUL + confidence",
                "Agent: risk + recommendation",
                "Dashboard: user + audit",
            ],
            "architecture",
            [
                "This is the system-level contribution in one line.",
                "The predictive layer estimates RUL, the agent layer turns that into deterministic decision support, and the dashboard exposes both operator and audit views.",
            ],
            ["doc/paper/main.tex", "out/dashboard_layer/01_ui_backend_contract_v1.json"],
        ),
        Slide(
            "Predictive Layer",
            "Best model \u2260 deployed model",
            [
                "Multiple models tested: RF, HGB, LSTM, GRU",
                "GRU achieved the best raw RMSE",
                "RF was deployed as champion",
                "Deployment needs stability, latency, and robustness",
            ],
            "tradeoff",
            [
                "This is the main predictive insight of the project.",
                "Even though GRU had the best benchmark score, RF was the better deployed choice because the system also had to satisfy stability and operational constraints.",
            ],
            ["out/predictive_layer/champion_record.json", "out/predictive_layer/04_metrics_global_by_model.csv", "out/predictive_layer/04_latency_by_model.csv"],
        ),
        Slide(
            "Agent Layer",
            "Prediction becomes decision support",
            [
                "Inputs: RUL, confidence, service status",
                "Outputs: risk level, risk score, recommendation",
                "Deterministic and auditable",
                "LLM optional, never controls decisions",
            ],
            "minimal",
            [
                "The agent layer is what turns model output into something operationally meaningful.",
                "Importantly, the logic is deterministic, and the LLM is only optional support for interpretation rather than the source of truth.",
            ],
            ["src/agent_layer/risk_engine.py", "src/agent_layer/orchestrator.py", "out/agent_layer/06_scenario_assistant_policy.txt"],
        ),
        Slide(
            "Dashboard",
            "The integrated user-facing system",
            [
                "Summary \u2192 decision",
                "Analysis \u2192 uncertainty",
                "Scenarios \u2192 what-if comparison",
                "Audit \u2192 traceability",
            ],
            "dashboard",
            [
                "The dashboard is where the three-layer architecture becomes visible to users.",
                "It exposes a fast decision view, a deeper explanation layer, deterministic scenario comparisons, and a technical audit path.",
            ],
            ["fig/dashboard/dashboard_v1_4.png", "fig/dashboard/dashboard_v1_5.png", "fig/dashboard/dashboard_v1_9.png", "fig/dashboard/dashboard_v1_8.png"],
        ),
        Slide(
            "Results",
            "Engineering constraints matter",
            [
                "GRU \u2192 best accuracy",
                "RF \u2192 best system behavior",
                "Best model \u2260 deployed model",
                "PHM quality is not RMSE alone",
            ],
            "tradeoff",
            [
                "This slide reinforces the central result from a system perspective.",
                "A deployable PHM system has to optimize for more than leaderboard performance, and that is why the final model choice differs from the best raw benchmark model.",
            ],
            ["out/predictive_layer/champion_record.json", "out/predictive_layer/04_metrics_global_by_model.csv"],
        ),
        Slide(
            "Scenarios",
            "Interpretability through what-if analysis",
            [
                "Example: \u0394RUL \u2248 0",
                "Risk unchanged",
                "Not all changes matter",
                "The system reveals sensitivity, not just prediction",
            ],
            "scenario",
            [
                "Scenario analysis is useful because it shows how the deployed model responds to controlled changes.",
                "In some cases the effect is very small, which tells us that the edited variables are locally insensitive or not central to the deployed model behavior.",
            ],
            ["out/agent_layer/05_scenario_examples.json", "data/test_input_reference.md"],
        ),
        Slide(
            "Conclusion",
            "A deployable PHM system, not just a model",
            [
                "Calibrated prediction",
                "Deterministic decisions",
                "Traceability and auditability",
                "Main lesson: PHM quality \u2260 RMSE alone",
            ],
            "minimal",
            [
                "The final takeaway is that this project combines scientific modeling with engineering deployment discipline.",
                "What matters is not only predicting RUL, but making that prediction usable, explainable, and auditable in a real PHM workflow.",
            ],
            ["doc/paper/main.tex", "doc/paper/04_conclusion.md"],
        ),
    ]


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.65), Inches(0.4), Inches(12), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT

    st = slide.shapes.add_textbox(Inches(0.68), Inches(0.95), Inches(11.5), Inches(0.35))
    tf = st.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.size = Pt(14)
    run.font.color.rgb = MUTED


def add_footer(slide, sources: list[str], idx: int) -> None:
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12), Inches(0.2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Slide {idx} | Sources: " + "; ".join(sources[:2])
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT
        p.space_after = Pt(9)


def add_emphasis_badge(slide, text: str, left: float, top: float, width: float, fill_color: RGBColor, font_color: RGBColor = RGBColor(255, 255, 255)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = font_color


def add_image(slide, path: Path, left: float, top: float, width: float, height: float | None = None):
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_ppt(slides: list[Slide], assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_title(slide, spec.title, spec.subtitle)

        if spec.visual_type == "title":
            add_bullets(slide, spec.bullets, 0.85, 1.8, 7.1, 2.8)
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(2.0), Inches(3.7), Inches(2.0))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PANEL
            shape.line.color.rgb = ACCENT
            shape.line.width = Pt(2)
            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "CMAPSS\nPHM App"
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = ACCENT
            add_emphasis_badge(slide, "4-MINUTE VERSION", 8.95, 4.35, 2.9, TEAL)
        elif spec.visual_type == "problem":
            add_bullets(slide, spec.bullets, 0.75, 1.65, 4.9, 3.8)
            add_image(slide, assets["problem"], 5.95, 1.55, 6.7)
            add_emphasis_badge(slide, "NOT sufficient for PHM deployment", 0.8, 5.65, 4.8, RED)
        elif spec.visual_type == "architecture":
            add_bullets(slide, spec.bullets, 0.75, 1.7, 4.8, 3.9)
            add_image(slide, assets["architecture"], 5.9, 1.5, 6.6)
            add_emphasis_badge(slide, "Predictive \u2192 Agent \u2192 Dashboard", 6.5, 5.75, 5.0, TEAL)
        elif spec.visual_type == "tradeoff":
            add_bullets(slide, spec.bullets, 0.75, 1.6, 4.65, 4.1)
            add_image(slide, assets["tradeoff"], 5.8, 1.55, 6.7)
            add_emphasis_badge(slide, "Best model \u2260 deployed model", 6.55, 5.75, 4.5, RED)
        elif spec.visual_type == "dashboard":
            add_bullets(slide, spec.bullets, 0.75, 1.55, 4.3, 3.9)
            add_image(slide, assets["dashboard"], 4.9, 1.35, 7.7)
        elif spec.visual_type == "scenario":
            add_bullets(slide, spec.bullets, 0.75, 1.7, 4.7, 4.0)
            add_image(slide, assets["scenario"], 5.8, 1.65, 6.5)
        else:
            add_bullets(slide, spec.bullets, 0.85, 1.7, 10.8, 4.8)

        add_footer(slide, spec.sources, idx)
        slide.notes_slide.notes_text_frame.text = " ".join(spec.notes)

    prs.save(OUTPUT_PPTX)


def build_notes_file(slides: list[Slide]) -> None:
    lines = ["# Speaker Notes - 4-minute Presentation", ""]
    for idx, slide in enumerate(slides, start=1):
        lines.append(f"## Slide {idx} - {slide.title}")
        for note in slide.notes:
            lines.append(f"- {note}")
        lines.append("")
    OUTPUT_NOTES.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    ensure_dirs()
    metrics = load_metrics()
    assets = build_assets(metrics)
    slides = slide_specs()
    build_ppt(slides, assets)
    build_notes_file(slides)
    print(OUTPUT_PPTX)
    print(OUTPUT_NOTES)


if __name__ == "__main__":
    main()
