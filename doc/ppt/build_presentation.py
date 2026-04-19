from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from textwrap import wrap

import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parent
PROJECT_ROOT = ROOT.parents[1]
ASSET_DIR = ROOT / "assets"
PPTX_PATH = ROOT / "final_presentation.pptx"
PDF_PATH = ROOT / "final_presentation.pdf"
SLIDE_INDEX_PATH = ROOT / "final_slide_index.csv"
CLAIM_CHECKLIST_PATH = ROOT / "final_claim_checklist.txt"
README_PATH = ROOT / "README.md"

BG = RGBColor(248, 243, 232)
ACCENT = RGBColor(136, 94, 36)
TEXT = RGBColor(28, 36, 52)
MUTED = RGBColor(93, 99, 109)
TEAL = RGBColor(21, 110, 104)
LIGHT_PANEL = RGBColor(255, 255, 255)


@dataclass
class SlideSpec:
    title: str
    subtitle: str | None
    bullets: list[str]
    image_paths: list[Path]
    notes: str
    sources: list[str]


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def load_inputs() -> dict:
    global_metrics = pd.read_csv(PROJECT_ROOT / "out/predictive_layer/04_metrics_global_by_model.csv")
    latency = pd.read_csv(PROJECT_ROOT / "out/predictive_layer/04_latency_by_model.csv")
    per_dataset = pd.read_csv(PROJECT_ROOT / "out/predictive_layer/04_metrics_by_dataset_by_model.csv")
    champion = json.loads((PROJECT_ROOT / "out/predictive_layer/champion_record.json").read_text(encoding="utf-8"))
    scenario = json.loads((PROJECT_ROOT / "out/agent_layer/05_scenario_examples.json").read_text(encoding="utf-8"))
    return {
        "global_metrics": global_metrics,
        "latency": latency,
        "per_dataset": per_dataset,
        "champion": champion,
        "scenario": scenario,
    }


def save_plot(fig: plt.Figure, path: Path) -> None:
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches="tight", facecolor="#F8F3E8")
    plt.close(fig)


def make_motivation_chart(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 4))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    ax.axis("off")
    labels = [
        ("Point prediction only", "Insufficient for maintenance decisions"),
        ("Need uncertainty + risk", "Operators need confidence and severity"),
        ("Need traceable system", "Deployed PHM requires auditable interaction"),
    ]
    xs = [0.12, 0.5, 0.88]
    for x, (title, subtitle) in zip(xs, labels):
        ax.text(
            x,
            0.62,
            title,
            ha="center",
            va="center",
            fontsize=18,
            fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.6", fc="#FFFFFF", ec="#885E24", lw=2),
            color="#1C2434",
        )
        ax.text(x, 0.3, subtitle, ha="center", va="center", fontsize=13, color="#3E4956")
    ax.annotate("", xy=(0.38, 0.62), xytext=(0.24, 0.62), arrowprops=dict(arrowstyle="->", lw=2.5, color="#156E68"))
    ax.annotate("", xy=(0.76, 0.62), xytext=(0.62, 0.62), arrowprops=dict(arrowstyle="->", lw=2.5, color="#156E68"))
    ax.text(0.5, 0.92, "Why benchmark RUL models are not enough", ha="center", fontsize=20, fontweight="bold", color="#1C2434")
    save_plot(fig, path)


def make_architecture_chart(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 4.5))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    ax.axis("off")
    blocks = [
        (0.18, "Dashboard Layer", "Summary | Analysis | Scenarios | Audit", "#F6E8D1"),
        (0.5, "Agent Layer", "Deterministic risk, recommendation, scenarios", "#D9F1EE"),
        (0.82, "Predictive Layer", "Champion + fallback inference + confidence", "#F6E8D1"),
    ]
    for x, title, subtitle, fill in blocks:
        ax.text(
            x,
            0.58,
            f"{title}\n{subtitle}",
            ha="center",
            va="center",
            fontsize=16,
            fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.7", fc=fill, ec="#885E24", lw=2),
            color="#1C2434",
        )
    ax.annotate("", xy=(0.37, 0.58), xytext=(0.29, 0.58), arrowprops=dict(arrowstyle="->", lw=2.5, color="#156E68"))
    ax.annotate("", xy=(0.69, 0.58), xytext=(0.61, 0.58), arrowprops=dict(arrowstyle="->", lw=2.5, color="#156E68"))
    ax.annotate("", xy=(0.61, 0.4), xytext=(0.69, 0.4), arrowprops=dict(arrowstyle="->", lw=2.0, color="#885E24"))
    ax.text(0.5, 0.83, "Contract-based Predictive -> Agent -> Dashboard architecture", ha="center", fontsize=20, fontweight="bold", color="#1C2434")
    ax.text(0.33, 0.67, "request payload", fontsize=12, color="#3E4956")
    ax.text(0.64, 0.67, "trimmed predictive payload", fontsize=12, color="#3E4956")
    ax.text(0.62, 0.3, "RUL + confidence + trace", fontsize=12, color="#3E4956")
    save_plot(fig, path)


def make_global_metrics_chart(df: pd.DataFrame, path: Path) -> None:
    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    order = ["rf", "hgb", "gru", "lstm"]
    plot_df = df.set_index("model_name").loc[order].reset_index()
    x = range(len(plot_df))
    ax.bar([i - 0.18 for i in x], plot_df["rmse"], width=0.35, color="#885E24", label="RMSE")
    ax.bar([i + 0.18 for i in x], plot_df["mae"], width=0.35, color="#156E68", label="MAE")
    ax.set_xticks(list(x))
    ax.set_xticklabels(plot_df["model_name"].str.upper())
    ax.set_ylabel("Error")
    ax.set_title("Global predictive performance")
    ax.legend(frameon=False)
    for idx, value in enumerate(plot_df["rmse"]):
        ax.text(idx - 0.18, value + 0.25, f"{value:.2f}", ha="center", fontsize=9, color="#1C2434")
    save_plot(fig, path)


def make_latency_chart(df: pd.DataFrame, path: Path) -> None:
    fig, ax = plt.subplots(figsize=(9, 4.8))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    plot_df = df.sort_values("ms_per_sample", ascending=False)
    ax.bar(plot_df["model_name"].str.upper(), plot_df["ms_per_sample"], color=["#885E24", "#A77834", "#156E68", "#258C83"])
    ax.set_yscale("log")
    ax.set_ylabel("ms/sample (log scale)")
    ax.set_title("Serving latency tradeoff")
    for idx, (_, row) in enumerate(plot_df.iterrows()):
        ax.text(idx, row["ms_per_sample"] * 1.1, f"{row['ms_per_sample']:.3f}", ha="center", fontsize=9, color="#1C2434")
    save_plot(fig, path)


def make_scenario_chart(scenario: dict, path: Path) -> None:
    baseline = scenario["baseline"]
    scenario_out = scenario["scenario"]
    fig, ax = plt.subplots(figsize=(9.5, 4.8))
    fig.patch.set_facecolor("#F8F3E8")
    ax.set_facecolor("#F8F3E8")
    metrics = ["RUL", "Risk Score"]
    baseline_vals = [baseline["rul_pred"], baseline["risk_score"]]
    scenario_vals = [scenario_out["rul_pred"], scenario_out["risk_score"]]
    x = range(len(metrics))
    ax.bar([i - 0.18 for i in x], baseline_vals, width=0.35, color="#885E24", label="Baseline")
    ax.bar([i + 0.18 for i in x], scenario_vals, width=0.35, color="#156E68", label="Scenario")
    ax.set_xticks(list(x))
    ax.set_xticklabels(metrics)
    ax.set_ylabel("Value")
    ax.set_title("Deterministic baseline vs scenario comparison")
    ax.legend(frameon=False)
    for i, val in enumerate(baseline_vals):
        ax.text(i - 0.18, val + (1.2 if i == 0 else 0.8), f"{val:.2f}", ha="center", fontsize=9, color="#1C2434")
    for i, val in enumerate(scenario_vals):
        ax.text(i + 0.18, val + (1.2 if i == 0 else 0.8), f"{val:.2f}", ha="center", fontsize=9, color="#1C2434")
    save_plot(fig, path)


def create_assets(inputs: dict) -> dict[str, Path]:
    assets = {
        "motivation": ASSET_DIR / "motivation_problem.png",
        "architecture": ASSET_DIR / "architecture_layers.png",
        "global_metrics": ASSET_DIR / "global_metrics.png",
        "latency": ASSET_DIR / "latency_tradeoff.png",
        "scenario": ASSET_DIR / "scenario_comparison.png",
    }
    make_motivation_chart(assets["motivation"])
    make_architecture_chart(assets["architecture"])
    make_global_metrics_chart(inputs["global_metrics"], assets["global_metrics"])
    make_latency_chart(inputs["latency"], assets["latency"])
    make_scenario_chart(inputs["scenario"], assets["scenario"])
    return assets


def build_slide_specs(inputs: dict, assets: dict[str, Path]) -> list[SlideSpec]:
    champion = inputs["champion"]
    global_df = inputs["global_metrics"]
    rf_rmse = float(global_df.loc[global_df["model_name"] == "rf", "rmse"].iloc[0])
    gru_rmse = float(global_df.loc[global_df["model_name"] == "gru", "rmse"].iloc[0])
    scenario = inputs["scenario"]
    delta_rul = scenario["comparison"]["delta_rul_pred"]
    baseline_risk = scenario["baseline"]["risk_score"]
    scenario_risk = scenario["scenario"]["risk_score"]

    dashboard_imgs = [
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_4.png",
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_5.png",
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_9.png",
        PROJECT_ROOT / "fig/dashboard/dashboard_v1_8.png",
    ]

    return [
        SlideSpec(
            "Traceable Predictive Maintenance for NASA C-MAPSS",
            "An end-to-end PHM application with calibrated RUL prediction, deterministic risk logic, and auditable scenarios",
            [
                "Benchmark: NASA C-MAPSS",
                "System: Predictive Layer -> Agent Layer -> Dashboard Layer",
                "Focus: deployable PHM decision support rather than benchmark scores alone",
            ],
            [],
            "Open by framing the project as both a PHM study and a working AI system.",
            ["doc/paper/main.tex", "README.md"],
        ),
        SlideSpec(
            "Motivation and Problem",
            "Why point RUL prediction alone is not enough",
            [
                "Maintenance decisions require uncertainty, severity, and actionability.",
                "Benchmark optimization does not automatically produce a usable PHM application.",
                "This project closes that gap with calibrated prediction, deterministic decision logic, and traceable interaction.",
            ],
            [assets["motivation"]],
            "Use this slide to motivate the move from model-only research to system-level PHM support.",
            ["doc/paper/main.tex", "doc/paper/04_discussion.md"],
        ),
        SlideSpec(
            "NASA C-MAPSS Context",
            "Run-to-failure trajectories across four benchmark subsets",
            [
                "FD001-FD004 vary in operating conditions and fault complexity.",
                "Each unit is a multivariate cycle trajectory with operational settings and sensor measurements.",
                "The task is to estimate Remaining Useful Life before failure.",
            ],
            [],
            "Summarize the benchmark and remind the audience this is a cycle-based degradation problem.",
            ["data/readme.txt", "doc/paper/main.tex"],
        ),
        SlideSpec(
            "Project Objectives and Contributions",
            "What this system adds beyond benchmark prediction",
            [
                "A calibrated predictive layer with frozen preprocessing and champion selection.",
                "A deterministic agent layer for risk, recommendation, and scenario comparison.",
                "A dashboard with Summary, Analysis, Scenarios, and Technical Audit workflows.",
                "A reproducible artifact chain linking data, code, models, contracts, and audit traces.",
            ],
            [],
            "Make the contribution split explicit: prediction, decision support, user interface, and traceability.",
            ["doc/paper/main.tex", "doc/paper/03_results.md"],
        ),
        SlideSpec(
            "System Architecture",
            "Contract-driven Predictive -> Agent -> Dashboard stack",
            [
                "The dashboard collects manual or CSV-driven requests.",
                "The agent layer trims payloads, computes deterministic risk and recommendations, and manages scenarios.",
                "The predictive layer serves the champion model with a calibrated confidence band and trace metadata.",
            ],
            [assets["architecture"]],
            "Explain the layer boundaries and stress that contracts and service states are explicit.",
            ["doc/paper/main.tex", "out/dashboard_layer/01_ui_backend_contract_v1.json"],
        ),
        SlideSpec(
            "Data and Preprocessing",
            "From raw C-MAPSS rows to model-ready features",
            [
                "The deployed predictive model uses 16 selected features after low-variance filtering.",
                "Training uses a training-only StandardScaler reused at serving time.",
                "The operational target is capped RUL with cap 125, while raw RUL remains available for audit and ablations.",
            ],
            [],
            "Keep this practical: feature selection, scaling, and target policy are what matter for the deployed system.",
            ["out/eda/05_preprocessing_config.json", "out/eda/04_target_definition.txt"],
        ),
        SlideSpec(
            "Predictive Layer and Champion Selection",
            f"RF deployed despite GRU best raw RMSE ({gru_rmse:.2f} vs RF {rf_rmse:.2f})",
            [
                "Candidate models: RF, HGB, LSTM, and GRU.",
                f"Champion: {champion['champion_model']} | Fallback: {champion['fallback_model']}.",
                f"Selection reason: {champion['selection_reason']}",
                "The deployed choice balances quality, stability spread, latency, and integration readiness.",
            ],
            [assets["global_metrics"], assets["latency"]],
            "The main lesson is that the best benchmark model was not automatically the best deployed system model.",
            ["out/predictive_layer/champion_record.json", "out/predictive_layer/04_metrics_global_by_model.csv", "out/predictive_layer/04_latency_by_model.csv"],
        ),
        SlideSpec(
            "Agent Layer and Deterministic Decision Logic",
            "How RUL becomes risk, recommendations, and auditable scenario behavior",
            [
                "Risk logic uses predicted RUL, confidence-band width, and service status.",
                "Recommendations are deterministic and aligned with the risk state.",
                "The LLM is optional and non-central: it may assist interpretation, but never controls execution or comparison values.",
            ],
            [],
            "Stress that the decision layer remains deterministic even when an LLM is available.",
            ["src/agent_layer/risk_engine.py", "src/agent_layer/orchestrator.py", "out/agent_layer/06_scenario_assistant_policy.txt"],
        ),
        SlideSpec(
            "Dashboard Layer and User Workflow",
            "The UI operationalizes prediction, explanation, scenarios, and auditability",
            [
                "Summary: decision-level RUL, risk, urgency, and action.",
                "Analysis: explanation, uncertainty, temporal behavior, and fleet context.",
                "Scenarios: deterministic baseline-vs-scenario comparison.",
                "Technical Audit: trace fields, model version, service state, and contracts.",
            ],
            dashboard_imgs,
            "Use the screenshots as evidence that the integrated system is implemented, not only described.",
            ["fig/dashboard/dashboard_v1_4.png", "fig/dashboard/dashboard_v1_5.png", "fig/dashboard/dashboard_v1_8.png", "fig/dashboard/dashboard_v1_9.png"],
        ),
        SlideSpec(
            "Experiments and Results",
            "Evidence combines prediction quality, stability, and engineering tradeoffs",
            [
                "Global metrics: GRU best raw RMSE; RF strongest stable tabular baseline.",
                "Per-dataset metrics show why stability spread matters for deployment.",
                "Latency and fallback behavior are treated as first-class engineering evidence.",
            ],
            [assets["global_metrics"], assets["latency"]],
            "Keep the framing evidence-based: show both performance and deployability tradeoffs.",
            ["out/predictive_layer/04_metrics_global_by_model.csv", "out/predictive_layer/04_metrics_by_dataset_by_model.csv", "out/predictive_layer/04_latency_by_model.csv"],
        ),
        SlideSpec(
            "Scenarios and Interpretability",
            "Deterministic what-if analysis with optional non-central LLM assistance",
            [
                "Scenario parsing is deterministic-first and only uses the LLM as a fallback for intent interpretation.",
                f"Example scenario: baseline vs scenario delta RUL = {delta_rul:.3f} cycles.",
                f"Risk score remained {baseline_risk} -> {scenario_risk}, showing how some edits have low local impact.",
                "This makes scenario outputs useful for comparative sensitivity analysis rather than speculative control.",
            ],
            [assets["scenario"], PROJECT_ROOT / "fig/dashboard/dashboard_v1_9.png"],
            "Explain that low-impact scenario outcomes are often a model-feature and local-sensitivity issue, not a parser failure.",
            ["out/agent_layer/05_scenario_examples.json", "data/test_input_reference.md"],
        ),
        SlideSpec(
            "Engineering Constraints",
            "Deployment decisions were shaped by cost, latency, safety, reliability, and traceability",
            [
                "Cost: the LLM remains optional and non-central.",
                "Latency: the deployed RF path is simple and fast to serve.",
                "Safety: validation rules and protected identifiers constrain scenario execution.",
                "Reliability: explicit fallback and service-state propagation are surfaced end-to-end.",
                "Traceability: audit records and contracts support debugging and review.",
            ],
            [],
            "This slide is where the engineering maturity of the project should become obvious.",
            ["out/agent_layer/08_failure_modes.txt", "out/predictive_layer/08_deploy_checklist.txt", "out/dashboard_layer/08_state_handling_guide.txt"],
        ),
        SlideSpec(
            "Limitations and Future Work",
            "What the current system does well and where it still needs extension",
            [
                "The deployed champion is a single-row tabular model rather than a temporal serving model.",
                "Some scenario edits have low local sensitivity because they target non-features or weak local directions.",
                "The confidence-band policy is useful but relatively simple.",
                "Future work: richer local explanation, broader scenario libraries, stronger fleet analytics, and live deployment experiments.",
            ],
            [],
            "Be explicit and honest here; it strengthens the credibility of the rest of the deck.",
            ["doc/paper/04_discussion.md", "out/predictive_layer/08_residual_risks.md"],
        ),
        SlideSpec(
            "Conclusion",
            "The project demonstrates a deployable and traceable PHM system, not only a benchmark model",
            [
                "Scientifically: calibrated RUL modeling with documented champion selection.",
                "Engineering-wise: deterministic decision support, scenarios, and traceability integrated into one user-facing system.",
                "Main lesson: deployment quality depends on more than raw RMSE.",
            ],
            [],
            "Close by joining the scientific and engineering contributions into one claim.",
            ["doc/paper/main.tex", "doc/paper/04_conclusion.md"],
        ),
        SlideSpec(
            "Backup / Appendix",
            "Extra evidence for Q&A",
            [
                "Full per-dataset metrics and latency table.",
                "Selected feature list and service-state taxonomy.",
                "Contract field mapping and reproducibility commands.",
                "Scenario parsing examples and failure-mode handling.",
            ],
            [],
            "Use this only when questions go deeper into reproducibility, traceability, or failure cases.",
            ["environment.yml", "USER_GUIDE.md", "out/dashboard_layer/03_mapping_fields_table.csv"],
        ),
    ]


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    p.alignment = align
    return box


def add_bullets(slide, bullets: list[str], left, top, width, height) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
    return box


def add_footer(slide, text: str) -> None:
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.25), text, font_size=10, color=MUTED)


def add_image(slide, path: Path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(str(path), left, top, width=width)
    else:
        slide.shapes.add_picture(str(path), left, top, height=height)


def create_pptx(slides: list[SlideSpec]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        add_textbox(slide, Inches(0.55), Inches(0.35), Inches(11.8), Inches(0.55), spec.title, font_size=28, bold=True)
        if spec.subtitle:
            add_textbox(slide, Inches(0.58), Inches(0.95), Inches(11.5), Inches(0.4), spec.subtitle, font_size=14, color=MUTED)

        if idx == 1:
            add_bullets(slide, spec.bullets, Inches(0.9), Inches(2.0), Inches(7.2), Inches(2.5))
            add_textbox(slide, Inches(8.65), Inches(2.0), Inches(3.7), Inches(2.2), "CMAPSS PHM\nPresentation Deck", font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(8.55), Inches(4.4), Inches(3.9), Inches(1.0), "Generated from project artifacts\nApril 2026", font_size=14, color=MUTED, align=PP_ALIGN.CENTER)
        elif idx in (2, 5, 7, 10, 11):
            add_bullets(slide, spec.bullets, Inches(0.7), Inches(1.5), Inches(5.1), Inches(4.8))
            if len(spec.image_paths) == 1:
                add_image(slide, spec.image_paths[0], Inches(6.0), Inches(1.45), width=Inches(6.7))
            elif len(spec.image_paths) >= 2:
                add_image(slide, spec.image_paths[0], Inches(5.9), Inches(1.4), width=Inches(3.1))
                add_image(slide, spec.image_paths[1], Inches(9.15), Inches(1.4), width=Inches(3.1))
        elif idx == 9:
            add_bullets(slide, spec.bullets, Inches(0.65), Inches(1.4), Inches(4.4), Inches(3.5))
            positions = [
                (Inches(5.1), Inches(1.35)),
                (Inches(9.1), Inches(1.35)),
                (Inches(5.1), Inches(4.15)),
                (Inches(9.1), Inches(4.15)),
            ]
            for pth, (left, top) in zip(spec.image_paths, positions):
                add_image(slide, pth, left, top, width=Inches(3.65))
        else:
            add_bullets(slide, spec.bullets, Inches(0.8), Inches(1.55), Inches(11.5), Inches(4.9))

        add_footer(slide, f"Slide {idx} | Sources: " + "; ".join(spec.sources[:2]))
        slide.notes_slide.notes_text_frame.text = spec.notes + "\nSources: " + "; ".join(spec.sources)

    prs.save(PPTX_PATH)


def draw_wrapped(c: canvas.Canvas, text: str, x: float, y: float, width_chars: int, font="Helvetica", size=20, leading=24):
    c.setFont(font, size)
    yy = y
    for line in wrap(text, width_chars):
        c.drawString(x, yy, line)
        yy -= leading
    return yy


def create_pdf(slides: list[SlideSpec]) -> None:
    page_size = landscape((13.333 * inch, 7.5 * inch))
    c = canvas.Canvas(str(PDF_PATH), pagesize=page_size)
    page_w, page_h = page_size

    for idx, spec in enumerate(slides, start=1):
        c.setFillColorRGB(248 / 255, 243 / 255, 232 / 255)
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        c.setFillColorRGB(28 / 255, 36 / 255, 52 / 255)
        c.setFont("Helvetica-Bold", 24)
        c.drawString(0.55 * inch, page_h - 0.6 * inch, spec.title)
        if spec.subtitle:
            c.setFont("Helvetica", 12)
            c.setFillColorRGB(93 / 255, 99 / 255, 109 / 255)
            c.drawString(0.58 * inch, page_h - 0.95 * inch, spec.subtitle)

        c.setFillColorRGB(28 / 255, 36 / 255, 52 / 255)
        y = page_h - 1.6 * inch
        x = 0.8 * inch
        for bullet in spec.bullets:
            c.setFont("Helvetica-Bold", 16)
            c.drawString(x, y, u"\u2022")
            y = draw_wrapped(c, bullet, x + 0.25 * inch, y, 62, size=15, leading=19) - 0.08 * inch

        image_slots = []
        if len(spec.image_paths) == 1:
            image_slots = [(spec.image_paths[0], 6.1 * inch, 1.5 * inch, 6.4 * inch, 4.4 * inch)]
        elif len(spec.image_paths) == 2:
            image_slots = [
                (spec.image_paths[0], 6.0 * inch, 3.5 * inch, 3.0 * inch, 2.1 * inch),
                (spec.image_paths[1], 9.2 * inch, 3.5 * inch, 3.0 * inch, 2.1 * inch),
            ]
        elif len(spec.image_paths) == 4:
            image_slots = [
                (spec.image_paths[0], 5.2 * inch, 3.6 * inch, 3.45 * inch, 2.0 * inch),
                (spec.image_paths[1], 8.9 * inch, 3.6 * inch, 3.45 * inch, 2.0 * inch),
                (spec.image_paths[2], 5.2 * inch, 1.1 * inch, 3.45 * inch, 2.0 * inch),
                (spec.image_paths[3], 8.9 * inch, 1.1 * inch, 3.45 * inch, 2.0 * inch),
            ]
        for img_path, left, bottom, width, height in image_slots:
            if img_path.exists():
                c.drawImage(str(img_path), left, bottom, width=width, height=height, preserveAspectRatio=True, anchor='c')

        c.setFillColorRGB(93 / 255, 99 / 255, 109 / 255)
        c.setFont("Helvetica", 9)
        c.drawString(0.55 * inch, 0.35 * inch, f"Slide {idx} | Sources: " + "; ".join(spec.sources[:2]))
        c.showPage()

    c.save()


def write_slide_index(slides: list[SlideSpec]) -> None:
    rows = [
        {
            "slide_no": idx,
            "title": spec.title,
            "subtitle": spec.subtitle or "",
            "primary_sources": "; ".join(spec.sources),
            "notes": spec.notes,
        }
        for idx, spec in enumerate(slides, start=1)
    ]
    pd.DataFrame(rows).to_csv(SLIDE_INDEX_PATH, index=False)


def write_claim_checklist(inputs: dict) -> None:
    champion = inputs["champion"]
    lines = [
        "Final slide claim checklist",
        "",
        "- Claim: the system is implemented as a Predictive -> Agent -> Dashboard stack.",
        "  Evidence: doc/paper/main.tex; src/predictive_layer/*; src/agent_layer/*; src/dashboard_layer/*",
        "",
        f"- Claim: the deployed champion is {champion['champion_model']} and fallback is {champion['fallback_model']}.",
        "  Evidence: out/predictive_layer/champion_record.json",
        "",
        f"- Claim: {champion['best_overall_model']} achieved the best raw RMSE but was not promoted.",
        "  Evidence: out/predictive_layer/champion_record.json; out/predictive_layer/04_metrics_global_by_model.csv",
        "",
        "- Claim: risk and recommendations are deterministic.",
        "  Evidence: src/agent_layer/risk_engine.py; src/agent_layer/recommender.py",
        "",
        "- Claim: scenario execution is deterministic and auditable even when LLM assistance is available.",
        "  Evidence: src/agent_layer/orchestrator.py; src/agent_layer/scenario_rules.py; out/agent_layer/06_scenario_assistant_policy.txt",
        "",
        "- Claim: the deck screenshots reflect the implemented dashboard tabs.",
        "  Evidence: fig/dashboard/dashboard_v1_4.png; fig/dashboard/dashboard_v1_5.png; fig/dashboard/dashboard_v1_8.png; fig/dashboard/dashboard_v1_9.png",
        "",
        "- Claim: every major slide in the deck is tied to a project artifact.",
        "  Evidence: doc/ppt/final_slide_index.csv; doc/ppt/02_evidence_inventory.csv",
    ]
    CLAIM_CHECKLIST_PATH.write_text("\n".join(lines), encoding="utf-8")


def write_readme() -> None:
    text = """# PPT package

This folder contains the final CMAPSS PHM presentation package.

## Main outputs
- `final_presentation.pptx`: presentation deck
- `final_presentation.pdf`: PDF export generated from the same slide content
- `final_slide_index.csv`: slide-to-source mapping
- `final_claim_checklist.txt`: claim validation checklist

## Support files
- `01_slide_outline.md`
- `01_storyline_map.csv`
- `02_evidence_inventory.csv`
- `02_slide_content.md`
- `assets/`

## Rebuild
Run:

```powershell
C:\\Users\\quant\\miniconda3\\envs\\cmapss\\python.exe doc\\ppt\\build_presentation.py
```
"""
    README_PATH.write_text(text, encoding="utf-8")


def main() -> None:
    ensure_dirs()
    inputs = load_inputs()
    assets = create_assets(inputs)
    slides = build_slide_specs(inputs, assets)
    create_pptx(slides)
    create_pdf(slides)
    write_slide_index(slides)
    write_claim_checklist(inputs)
    write_readme()
    print(PPTX_PATH)
    print(PDF_PATH)


if __name__ == "__main__":
    main()
